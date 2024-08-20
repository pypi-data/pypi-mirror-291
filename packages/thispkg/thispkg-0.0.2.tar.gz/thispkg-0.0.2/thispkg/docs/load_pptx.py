from typing import Generator

from pptx import Presentation

from ..base64 import base64
from ._base import RawDoc


class PptxLoader(RawDoc):
    def extract_text(self) -> Generator[str, None, None]:
        with self.get_file() as file:
            prs = Presentation(file)
            for slide in prs.slides:  # type: ignore
                for shape in slide.shapes:  # type: ignore
                    if shape.has_text_frame:  # type: ignore
                        text_frame = shape.text_frame  # type: ignore
                        for paragraph in text_frame.paragraphs:  # type: ignore
                            if paragraph.text:  # type: ignore
                                yield paragraph.text  # type: ignore
                            else:
                                continue

    def extract_image(self):
        with self.get_file() as file:
            prs = Presentation(file)
            for slide in prs.slides:  # type: ignore
                for shape in slide.shapes:  # type: ignore
                    if shape.shape_type == 13:  # type: ignore
                        image = shape.image  # type: ignore
                        yield base64.b64encode(image.blob).decode()  # type: ignore
                    else:
                        continue
